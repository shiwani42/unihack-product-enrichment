#!/usr/bin/env python3
"""Fill the official UniHack prototype template. Always rebuilds from the pristine file.

Usage: PYTHONPATH=. python3 scripts/build_deck.py
"""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "guidelines" / "[EXT] UniHack-Protoype Template .pptx"
SHOTS = ROOT / "demo_build" / "deck_shots"
FALLBACK = ROOT / "demo_build" / "screenshots"
CROP = ROOT / "demo_build" / "deck_crops"
OUT = ROOT / "submission" / "UniHack_thExplorers_Prototype.pptx"

INK = RGBColor(0x11, 0x18, 0x27)
MUTED = RGBColor(0x4B, 0x55, 0x63)
GREEN = RGBColor(0x06, 0x76, 0x47)
MINT = RGBColor(0xEC, 0xF8, 0xF2)
HAIR = RGBColor(0xE5, 0xE7, 0xEB)
PALE = RGBColor(0xF8, 0xF9, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NAVY = RGBColor(0x00, 0x38, 0x78)
SANS = "Calibri"
MONO = "Consolas"

GITHUB = "https://github.com/shiwani42/unihack-product-enrichment"
DEMO = "https://vimeo.com/1220615209"
PROTO = "https://unilog-tau.vercel.app"

# Official template: 10.00 x 5.625. Navy header ~0.61", cyan footer from ~5.51".
L, R = 0.42, 9.58
W = R - L
TOP = 1.66
BOT = 5.34
H = BOT - TOP


def box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    return tb, tf


def style(par, text, size=12, bold=False, color=INK, font=SANS, align=None,
          space_after=3, space_before=0):
    par.clear()
    par.alignment = align if align is not None else PP_ALIGN.LEFT
    par.space_before = Pt(space_before)
    par.space_after = Pt(space_after)
    run = par.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return par


def fill_tf(tf, items):
    first = True
    for text, kw in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        style(p, text, **kw)


def card(slide, x, y, w, h, fill=WHITE, border=HAIR, radius=0.08):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border
    sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def rect(slide, x, y, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def wipe_text(slide, keep_titles=True):
    """Remove leftover instructional text. Keep pictures and official titles."""
    titles = ("guidelines", "team details", "brief about", "opportunities",
              "list of features", "process flow", "wireframes", "architecture",
              "technologies", "estimated implementation", "snapshots",
              "additional details", "provide links", "thank")
    for shape in list(slide.shapes):
        if not shape.has_text_frame:
            continue
        if keep_titles:
            t = (shape.text_frame.text or "").strip().lower()
            if any(t.startswith(k) for k in titles) and len(t) < 90:
                continue
            if shape.is_placeholder and "how does your solution" in t:
                slide.shapes._spTree.remove(shape._element)
                continue
        if shape.is_placeholder or True:
            # drop long instructional copy; keep short official titles
            t = (shape.text_frame.text or "").strip()
            if keep_titles and len(t) < 90 and "\n" not in t:
                continue
            if keep_titles and t.lower().startswith(("guidelines", "team details",
                                                     "brief about your solution",
                                                     "opportunities",
                                                     "list of features offered",
                                                     "process flow diagram",
                                                     "wireframes",
                                                     "architecture diagram",
                                                     "technologies used",
                                                     "estimated implementation",
                                                     "snapshots of the mvp",
                                                     "additional details",
                                                     "provide links to your")):
                continue
            slide.shapes._spTree.remove(shape._element)


def crop_png(src: Path, dest: Path, box_px):
    dest.parent.mkdir(parents=True, exist_ok=True)
    with Image.open(src) as im:
        im.convert("RGB").crop(box_px).save(dest, "PNG")
    return dest


def prepare_crops():
    CROP.mkdir(parents=True, exist_ok=True)
    hero = SHOTS / "hero.png"
    rec = SHOTS / "record_panel.png"
    rec_wp = SHOTS / "record_panel_wp.png"
    drawer_e = SHOTS / "drawer_evidence.png"
    drawer_r = SHOTS / "drawer_record.png"
    catalog = SHOTS / "catalog_page.png"
    batch = SHOTS / "batch_progress.png"

    # hero: workbench only (drop giant headline so type stays readable)
    if hero.exists():
        with Image.open(hero) as im:
            w, h = im.size
            crop_png(hero, CROP / "workbench.png", (0, int(h * 0.42), w, h))
    if rec.exists():
        with Image.open(rec) as im:
            w, h = im.size
            crop_png(rec, CROP / "record.png", (0, 0, w, min(h, int(w * 0.92))))
    if rec_wp.exists():
        with Image.open(rec_wp) as im:
            w, h = im.size
            crop_png(rec_wp, CROP / "record_wp.png", (0, 0, w, min(h, int(w * 0.92))))
    if drawer_e.exists():
        with Image.open(drawer_e) as im:
            w, h = im.size
            crop_png(drawer_e, CROP / "evidence.png", (0, 0, w, min(h, int(w * 1.05))))
    if drawer_r.exists():
        with Image.open(drawer_r) as im:
            w, h = im.size
            crop_png(drawer_r, CROP / "drawer_record.png", (0, 0, w, min(h, int(w * 1.05))))
    if catalog.exists():
        with Image.open(catalog) as im:
            w, _h = im.size
            crop_png(catalog, CROP / "catalog.png", (0, 0, w, int(w * 0.48)))
    if batch.exists():
        crop_png(batch, CROP / "batch.png", (0, 0, *Image.open(batch).size))

    # fallbacks from older screenshots if a crop is missing
    for name, src in (
        ("workbench.png", FALLBACK / "hero.png"),
        ("record.png", FALLBACK / "enrich_result.png"),
        ("evidence.png", FALLBACK / "drawer_evidence.png"),
        ("catalog.png", FALLBACK / "catalog_table.png"),
    ):
        if not (CROP / name).exists() and src.exists():
            crop_png(src, CROP / name, (0, 0, *Image.open(src).size))


def framed(slide, path, x, y, w, max_h=None, caption=None):
    """Place an image at width w, never exceeding max_h. Thin hairline frame."""
    path = Path(path)
    with Image.open(path) as im:
        ar = im.height / im.width
    h = w * ar
    if max_h and h > max_h:
        h = max_h
        w = h / ar
    img = slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                   width=Inches(w), height=Inches(h))
    fr = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    fr.fill.background()
    fr.line.color.rgb = HAIR
    fr.line.width = Pt(0.75)
    fr.shadow.inherit = False
    # send frame behind picture
    tree = slide.shapes._spTree
    tree.remove(fr._element)
    tree.insert(list(tree).index(img._element), fr._element)
    used_h = h
    if caption:
        _, tf = box(slide, x, y + h + 0.04, w, 0.22)
        fill_tf(tf, [(caption, dict(size=9, color=MUTED, space_after=0))])
        used_h += 0.26
    return w, used_h


def stat_card(slide, x, y, w, h, value, label, fill=PALE):
    card(slide, x, y, w, h, fill=fill, border=HAIR)
    _, tf = box(slide, x + 0.10, y + 0.10, w - 0.20, h - 0.18)
    fill_tf(tf, [
        (value, dict(size=20, bold=True, color=GREEN, space_after=1)),
        (label, dict(size=10, color=MUTED, space_after=0)),
    ])


def main() -> None:
    prepare_crops()
    prs = Presentation(str(TEMPLATE))
    slides = list(prs.slides)

    # ========== 0 Title (white canvas under navy header) ==========
    s = slides[0]
    for shape in list(s.shapes):
        if shape.has_text_frame:
            s.shapes._spTree.remove(shape._element)
    _, tf = box(s, 0.50, 0.95, 9.0, 2.20)
    fill_tf(tf, [
        ("unilog enrichment engine", dict(size=32, bold=True, color=NAVY, space_after=8)),
        ("Evidence-first product intelligence for industrial commerce.",
         dict(size=16, color=INK, space_after=6)),
        ("Distributor input in. Catalog row out. Every value traceable to its source.",
         dict(size=14, bold=True, color=GREEN, space_after=4)),
        ("Team thExplorers  ·  Shiwani Mishra  ·  Saurabh Gupta  ·  UniHack 2026",
         dict(size=12, color=MUTED, space_after=0)),
    ])
    metrics = [
        ("100%", "field match vs organizer\nexpected output (134 / 134)"),
        ("1,000 / 1,000", "input rows classified\nto a leaf template"),
        ("77 tests", "hermetic suite, ~2 s\noffline by default"),
        ("$0.0004", "compute per SKU\non the rules path"),
    ]
    for i, (val, lab) in enumerate(metrics):
        x = 0.50 + i * 2.35
        card(s, x, 3.48, 2.22, 1.28, fill=PALE)
        _, tf = box(s, x + 0.12, 3.56, 1.98, 1.12)
        fill_tf(tf, [
            (val, dict(size=18, bold=True, color=GREEN, space_after=2)),
            (lab, dict(size=10, color=MUTED, space_after=0)),
        ])

    # ========== 1 Team — wipe ALL leftover template copy ==========
    s = slides[1]
    for shape in list(s.shapes):
        if shape.has_text_frame:
            s.shapes._spTree.remove(shape._element)
    _, tf = box(s, 0.55, 3.58, 8.9, 1.50)
    fill_tf(tf, [
        ("Team details", dict(size=14, bold=True, color=NAVY, space_after=8)),
        ("Team name:   thExplorers", dict(size=16, bold=True, space_after=4)),
        ("Team leader:   Shiwani Mishra", dict(size=15, space_after=3)),
        ("Member:   Saurabh Gupta", dict(size=15, space_after=0)),
    ])

    # ========== 2 Brief ==========
    s = slides[2]
    wipe_text(s)
    _, tf = box(s, L, TOP, W, 1.15)
    fill_tf(tf, [
        ("Distributors hand us an input table. We return a Unilog "
         "catalog row. Every filled cell is traced to a source URL.",
         dict(size=15, bold=True, space_after=6)),
        ("Manufacturer-first sourcing (Amazon / eBay blocked) → identity → "
         "leaf classification → HTML / JSON-LD / PDF extraction → unit & LOV "
         "normalisation → five governed descriptions → validation. Rules do the "
         "work. An LLM is last-mile only, off by default. Blank beats invented.",
         dict(size=12, color=MUTED, space_after=0)),
    ])
    stats = [
        ("100%", "Reference field match\nPDSH4816AF · WDTS7024RZ"),
        ("39.28", "Avg fields filled\nacross 1,000 rows"),
        ("26", "Manufacturer domains\ncited in provenance"),
        ("$0.0004", "Rules-path cost\nper SKU, zero APIs"),
    ]
    for i, (val, lab) in enumerate(stats):
        x = L + i * (W + 0.12) / 4
        stat_card(s, x, 3.00, 2.20, 1.18, val, lab, fill=MINT if i == 0 else PALE)

    # ========== 3 Three questions ==========
    s = slides[3]
    wipe_text(s, keep_titles=False)
    # restore a title since we wiped the placeholder
    _, tf = box(s, L, 0.86, W, 0.42)
    fill_tf(tf, [("How it enriches, stays accurate, and scales",
                  dict(size=20, bold=True, space_after=0))])
    qa = [
        ("01  Enrichment",
         "Input analysis → de-dup merge → brand / MPN identity → leaf routing "
         "across 14 templates → manufacturer HTML, JSON-LD and PDF extraction "
         "→ unit & LOV normalisation → five descriptions → 252-column delivery."),
        ("02  Accuracy & trust",
         "100% vs the organizer expected output (63/63 and 71/71). Per-cell "
         "source URLs. High confidence requires external manufacturer evidence. "
         "Validators: LOV, char limits, UOM style, attribute sanity. Integrity "
         "tests ban fabricated defaults."),
        ("03  Scale",
         "Stateless row pipeline, parallel workers, cache-first fetch under a "
         "hard network budget. SSE for live ops, CLI for millions of rows. "
         "Add a category = one JSON template. 1,000 rows offline in ~60 s."),
    ]
    cw = (W - 0.24) / 3
    for i, (title, body) in enumerate(qa):
        x = L + i * (cw + 0.12)
        card(s, x, TOP, cw, 3.52, fill=PALE)
        _, tf = box(s, x + 0.14, TOP + 0.12, cw - 0.28, 3.28)
        fill_tf(tf, [
            (title, dict(size=13, bold=True, color=GREEN, space_after=8)),
            (body, dict(size=11, color=INK, space_after=0)),
        ])

    # ========== 4 Opportunities / USP ==========
    s = slides[4]
    wipe_text(s)
    cols = [
        ("Different",
         "Most enrichment tools fill fields. Ours makes every filled cell "
         "auditable: a provenance drawer opens the exact source URL, and "
         "confidence bands say what a buyer can trust without re-checking."),
        ("USP",
         "Audit-proof enrichment — every value verifiable in one click. "
         "Accuracy is scored field-by-field against the organizer’s own "
         "expected output, not claimed."),
        ("Fit",
         "Manual work drops from research-per-SKU to review-by-exception. "
         "Rows without manufacturer evidence are flagged, never guessed. "
         "That is the trust bar a 100% accuracy target requires."),
    ]
    cw = (W - 0.24) / 3
    for i, (title, body) in enumerate(cols):
        x = L + i * (cw + 0.12)
        card(s, x, TOP, cw, 3.52, fill=WHITE)
        accent = rect(s, x, TOP, 0.08, 3.52, GREEN)
        _, tf = box(s, x + 0.22, TOP + 0.16, cw - 0.36, 3.20)
        fill_tf(tf, [
            (title, dict(size=14, bold=True, space_after=8)),
            (body, dict(size=12, color=MUTED, space_after=0)),
        ])

    # ========== 5 Features ==========
    s = slides[5]
    wipe_text(s)
    left = [
        "Distributor input to catalog CSV / XLSX, headers untouched",
        "Manufacturer-first sourcing; marketplace blocklist",
        "14 leaf templates + generic industrial fallback",
        "Attribute slots filled only when evidence exists",
        "Per-value provenance JSON — source URL per cell",
        "Five description types inside character limits",
    ]
    right = [
        "Marketing copy & features only from manufacturer pages",
        "Honest Actual Image flag — Yes only with mfr imagery",
        "Validators: LOV, limits, UOM, ecommerce block, sanity",
        "Confidence bands: high / medium / review, evidence-gated",
        "Live UI: enrich, catalog, evidence drawer, SSE batch",
        "Reference harness + 77 hermetic tests on every change",
    ]
    card(s, L, TOP, 4.46, 3.52, fill=PALE)
    card(s, L + 4.70, TOP, 4.46, 3.52, fill=PALE)
    _, tf = box(s, L + 0.18, TOP + 0.16, 4.10, 3.20)
    fill_tf(tf, [(f"•  {t}", dict(size=12, space_after=8)) for t in left])
    _, tf = box(s, L + 4.88, TOP + 0.16, 4.10, 3.20)
    fill_tf(tf, [(f"•  {t}", dict(size=12, space_after=8)) for t in right])

    # ========== 6 Process flow ==========
    s = slides[6]
    wipe_text(s)
    stages = [
        ("1  Ingest", "input table, placeholders,\ndedupe merge, never drop"),
        ("2  Identity", "Brand aliases, DIB / E1,\nMPN prefix rules"),
        ("3  Classify", "Leaf routing across\n14 category templates"),
        ("4  Extract", "HTML · JSON-LD · PDF\ncache-first fetch"),
        ("5  Normalise", "Units, LOV, canonical\nbrand casing"),
        ("6  Compose", "5 governed descriptions\n+ delivery asset names"),
        ("7  Validate", "Rules, confidence bands,\nintegrity tests"),
        ("8  Deliver", "CSV / XLSX +\nper-cell provenance JSON"),
    ]
    cw, ch = 2.16, 1.28
    gap_x, gap_y = 0.16, 0.18
    for i, (title, sub) in enumerate(stages):
        r, c = divmod(i, 4)
        x = L + c * (cw + gap_x)
        y = TOP + r * (ch + gap_y)
        fill = MINT if i in (0, 7) else PALE
        card(s, x, y, cw, ch, fill=fill)
        _, tf = box(s, x + 0.12, y + 0.12, cw - 0.24, ch - 0.20)
        fill_tf(tf, [
            (title, dict(size=12, bold=True, color=GREEN if i in (0, 7) else INK,
                         space_after=4)),
            (sub, dict(size=10, color=MUTED, space_after=0)),
        ])
    _, tf = box(s, L, TOP + 2 * ch + gap_y + 0.10, W, 0.28)
    fill_tf(tf, [("Fail-safe: every row is wrapped — a bad SKU never kills the batch.",
                  dict(size=11, color=MUTED, space_after=0))])

    # ========== 7 Wireframes / product ==========
    s = slides[7]
    wipe_text(s)
    left_p = CROP / "workbench.png"
    right_p = CROP / "record.png"
    if left_p.exists():
        framed(s, left_p, L, TOP, 4.46, max_h=3.40,
               caption="Enrich workbench, input table")
    if right_p.exists():
        framed(s, right_p, L + 4.70, TOP, 4.46, max_h=3.40,
               caption="Result — Frigidaire PDSH4816AF, 100% vs reference")

    # ========== 8 Architecture ==========
    s = slides[8]
    wipe_text(s)
    layers = [
        ("Interface", "FastAPI · SSE live stream · static web UI · CLI"),
        ("Orchestration", "pipeline.enrich_input_row — fail-safe per row · dedup merge · workers"),
        ("Intelligence", "identity/ resolver · classify/ templates · extract/ HTML-PDF-JSON-LD · optional LLM"),
        ("Trust", "validate/ LOV + limits + sanity · provenance map · confidence bands · reference harness"),
        ("Delivery", "CSV · XLSX · provenance JSON · batch reports"),
    ]
    lh = 0.54
    for i, (name, tech) in enumerate(layers):
        y = TOP + i * (lh + 0.08)
        card(s, L, y, W, lh, fill=MINT if i == 3 else PALE)
        _, tf = box(s, L + 0.16, y + 0.10, 1.70, 0.36)
        fill_tf(tf, [(name, dict(size=13, bold=True, space_after=0))])
        _, tf = box(s, L + 1.95, y + 0.10, W - 2.20, 0.36)
        fill_tf(tf, [(tech, dict(size=12, color=MUTED, space_after=0))])
    _, tf = box(s, L, TOP + 5 * (lh + 0.08) + 0.02, W, 0.30)
    fill_tf(tf, [("Uniform cache-first fetch: retry / backoff, atomic writes, TTL, "
                  "UNILOG_FETCH_BUDGET. Playwright only on 403. UNILOG_LIVE_FETCH=0 kills the network.",
                  dict(size=10, color=MUTED, space_after=0))])

    # ========== 9 Technologies ==========
    s = slides[9]
    wipe_text(s)
    techs = [
        ("Python 3.11", "Pipeline core, typed modules"),
        ("FastAPI + SSE", "Upload, live stream, downloads"),
        ("httpx + Playwright", "Resilient fetch, 403 fallback"),
        ("BeautifulSoup + extruct", "HTML tables, JSON-LD, microdata"),
        ("pdfplumber", "Manufacturer PDF datasheets"),
        ("openpyxl", "Delivery-format XLSX"),
        ("pytest × 77", "Hermetic suite, offline default"),
        ("Remotion + FFmpeg", "Reproducible 3-minute demo film"),
    ]
    cw, ch = 4.46, 0.72
    for i, (name, why) in enumerate(techs):
        r, c = divmod(i, 2)
        x = L + c * (cw + 0.24)
        y = TOP + r * (ch + 0.10)
        card(s, x, y, cw, ch, fill=PALE)
        _, tf = box(s, x + 0.16, y + 0.10, cw - 0.32, 0.54)
        fill_tf(tf, [
            (name, dict(size=13, bold=True, font=MONO, space_after=1)),
            (why, dict(size=11, color=MUTED, space_after=0)),
        ])

    # ========== 10 Cost ==========
    s = slides[10]
    wipe_text(s)
    rows = [
        ("Approach", "Cost / SKU", "Notes", True),
        ("Manual enrichment (offshore)", "$5.00 – $15.00", "Industry rate for research + entry", False),
        ("unilog — rules path", "$0.0004", "CPU only · 1,000 rows ≈ 60 s · zero API calls", False),
        ("unilog — LLM last-mile (worst)", "$0.0204 max", "Hard-capped calls / run · off by default", False),
    ]
    rh = 0.52
    for i, (a, b, c, head) in enumerate(rows):
        y = TOP + i * (rh + 0.08)
        card(s, L, y, W, rh, fill=HAIR if head else WHITE)
        _, tfa = box(s, L + 0.16, y + 0.12, 3.70, 0.30)
        fill_tf(tfa, [(a, dict(size=12, bold=head, space_after=0))])
        _, tfb = box(s, L + 4.00, y + 0.12, 1.80, 0.30)
        fill_tf(tfb, [(b, dict(size=12, bold=True, color=INK if head else GREEN,
                              font=MONO, space_after=0))])
        _, tfc = box(s, L + 5.90, y + 0.12, 3.10, 0.30)
        fill_tf(tfc, [(c, dict(size=11, color=MUTED, space_after=0))])
    _, tf = box(s, L, TOP + 4 * (rh + 0.08) + 0.06, W, 0.50)
    fill_tf(tf, [("A 1M-SKU catalog: ~$400 compute on the rules path vs $5–15M "
                  "manually — four orders of magnitude, with provenance a human can audit.",
                  dict(size=13, bold=True, space_after=0))])

    # ========== 11 MVP snapshots — TWO images only, no overflow ==========
    s = slides[11]
    wipe_text(s)
    cat = CROP / "catalog.png"
    ev = CROP / "evidence.png"
    if cat.exists():
        framed(s, cat, L, TOP, 4.46, max_h=3.40,
               caption="Catalog — 1,000 enriched rows, ready for PIM")
    if ev.exists():
        framed(s, ev, L + 4.70, TOP, 4.46, max_h=3.40,
               caption="Evidence drawer — blank beats invented")

    # ========== 12 Future ==========
    s = slides[12]
    wipe_text(s)
    fut = [
        ("Organizer reference packs",
         "Drop-in importer already ships. Real LOV (~161k), UOM standards, "
         "27k brand list and the 200-row ground-truth scorer activate automatically."),
        ("Taxonomy at 14k scale",
         "Embedding matcher over the full leaf index, sitting beside the current rule router."),
        ("Human review queue",
         "Review-band rows land in a triage UI; approvers publish into a CX1-style PIM."),
        ("HyperScale agent fit",
         "Packaged as a merchandising-agent accelerator: same provenance contract, "
         "streaming connector for distributor ERPs."),
    ]
    cw, ch = 4.46, 1.64
    for i, (name, desc) in enumerate(fut):
        r, c = divmod(i, 2)
        x = L + c * (cw + 0.24)
        y = TOP + r * (ch + 0.16)
        card(s, x, y, cw, ch, fill=PALE)
        _, tf = box(s, x + 0.16, y + 0.14, cw - 0.32, ch - 0.24)
        fill_tf(tf, [
            (name, dict(size=13, bold=True, space_after=6)),
            (desc, dict(size=12, color=MUTED, space_after=0)),
        ])

    # ========== 13 Links ==========
    s = slides[13]
    wipe_text(s)
    links = [
        ("GitHub public repository", GITHUB),
        ("Demo video — 3 minutes", DEMO),
        ("Working prototype", PROTO),
    ]
    lh = 0.92
    for i, (label, url) in enumerate(links):
        y = TOP + i * (lh + 0.16)
        card(s, L, y, W, lh, fill=PALE)
        _, tf = box(s, L + 0.28, y + 0.14, W - 0.56, 0.68)
        fill_tf(tf, [
            (label, dict(size=13, bold=True, space_after=2)),
            (url, dict(size=14, font=MONO, color=GREEN, space_after=0)),
        ])

    # ========== 14 Closing — template already has Thank You baked in ==========
    s = slides[14]
    for shape in list(s.shapes):
        if shape.has_text_frame:
            s.shapes._spTree.remove(shape._element)
    _, tf = box(s, 0.50, 4.85, 6.4, 0.38)
    fill_tf(tf, [
        ("thExplorers  ·  Shiwani Mishra  ·  Saurabh Gupta",
         dict(size=12, color=WHITE, space_after=0)),
    ])

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"wrote {OUT}")


if __name__ == "__main__":
    main()
